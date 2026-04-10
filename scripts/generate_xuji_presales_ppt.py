from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

slides = [
    ("屯象OS × 徐记海鲜\n23套系统替换商务售前方案", ["对标天财商龙SaaS全模块，并以AI-Native经营中枢实现超越", "屯象科技｜2026"]),
    ("客户现状与痛点", ["徐记海鲜当前 23 套系统并行，数据割裂、流程断点、决策滞后", "门店链路复杂：预订/桌台/点单/收银/KDS/外卖", "总部管控复杂：审批、对账、报表、跨部门协同", "目标：12个月完成核心替换并沉淀标杆模板"]),
    ("方案目标：超越天财，不是复制天财", ["目标一：表意等价（天财验收语言全覆盖）", "目标二：架构代差（边缘+云、Ontology、Agent协同）", "目标三：经营升级（从看报表到系统给建议）"]),
    ("屯象OS已开发技术底座（代码级）", ["五层架构：L4多端前台 → L3 Agent OS → L2业务中台 → L1 Ontology → L0设备适配", "13个微服务 + 多端壳层（Android/Windows/iPad/KDS/PWA/小程序）", "PostgreSQL RLS多租户 + 事件总线 + Mac mini边缘同步"]),
    ("模块全景（对应天财能力域）", ["交易履约：tx-trade + web-pos + web-kds + web-reception", "商品菜单：tx-menu（菜品/BOM/定价/套餐/发布）", "会员增长：tx-member + tx-growth", "供应链：tx-supply；财务：tx-finance；组织运营：tx-org+tx-ops", "分析决策：tx-analytics + tx-intel + tx-agent"]),
    ("徐记23套系统替换总览", ["🟢 可替代 15 套 ｜ 🟡 基本可替代 6 套 ｜ 🟠 部分可替代 1 套 ｜ 🔴 不可替代 0 套", "结论：21/23（91%）已可替换或基本替换"]),
    ("23套系统替换映射（示例）", ["正品贵德POS → tx-trade + web-pos + web-kds", "微生活 → tx-member + tx-agent(private_ops)", "奥琦玮G10 → tx-supply", "金蝶 → tx-finance（保留凭证接口）", "红海云HR → tx-org；蓝凌OA → tx-ops"]),
    ("对标天财全模块：双轨策略", ["轨道A（必须覆盖）：POS/KDS/报表SKU/供应链台账/总部主数据", "轨道B（必须超越）：AI建议、跨域溯源、边缘韧性、决策审计", "商务表达：既能替代，又能持续增收降本"]),
    ("POS与门店交易域对比", ["已具备：点单、加单、收银、桌台、预订、外卖、KDS、打印桥接", "重点补齐：存酒、押金专模、协议挂账、快捷键、打印管理器可视化", "价值：门店效率提升 + 高风险操作可审计"]),
    ("总部报表域对比（296张 vs 屯象策略）", ["P0：固化高频经营+财务+执行报表SKU", "P1：叠加经营叙事引擎（自动摘要）", "P2：自然语言问数 + 角色化驾驶舱", "从“报表数量竞争”升级为“决策速度竞争”"]),
    ("超越点1：9大Agent + 73 Actions", ["覆盖折扣守护、排菜、出餐、会员、库存、财务、巡店、客服、私域", "4时间点推送（晨/午/战前/晚）", "传统SaaS记录结果；屯象OS提供可执行建议"]),
    ("超越点2：Ontology全链路经营底座", ["一笔业务可穿透‘人货场财组织’全链路", "问题可追因、决策可复盘、规则可沉淀", "降低总部与门店口径争议，统一经营语言"]),
    ("超越点3：边缘智能与稳定交付", ["安卓/Windows负责外设，Mac mini负责本地库+同步+边缘推理", "断网可运行、恢复可同步，适配高峰场景", "弱网门店稳定性优于纯云SaaS方案"]),
    ("23套系统组合方案（按经营中心）", ["交易履约中心｜商品与供应链中心｜客户经营中心", "运营与组织中心｜财务与结算中心｜总部决策中心"]),
    ("实施路线图（四阶段）", ["Phase1：先替换4套（POS/巡店/BI/消息推送）", "Phase2：扩展5套（菜单/会员/供应链/平台/iPad）", "Phase3：复杂域收口（财务/HR/OA）", "Phase4：旧系统归零（适配器逐步拔除）"]),
    ("里程碑与验收标准", ["M1总部可视：驾驶舱+关键报表", "M2门店可跑：交易高峰稳定", "M3会员可用：迁移+营销闭环", "M4供应链可控：采购/库存/损耗闭环"]),
    ("ROI与商务价值", ["直接价值：减少软件采购与运维成本、降低人工对账成本", "间接价值：缩短决策时滞、提升翻台与履约效率", "形成可复制的集团数字化经营模板"]),
    ("风险与保障", ["风险：网络、数据质量、人员迁移成本", "保障：双轨灰度、回退预案、分层培训、周度联合评审"]),
    ("商务打包建议", ["基础包：交易+会员+驾驶舱", "增强包：供应链+财务+组织绩效", "旗舰包：Agent经营助手+边缘AI+深度集成", "建议签约：订阅费 + 实施费 + 年度运营陪跑"]),
    ("结论与下一步", ["屯象OS已具备徐记23套系统替换主能力", "可完成对天财全模块覆盖，并以AI-Native实现超越", "建议立即启动‘样板店-区域-全集团’三段式落地"]),
]


def add_slide(prs, title, bullets):
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title

    title_tf = slide.shapes.title.text_frame
    title_tf.paragraphs[0].font.size = Pt(34)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = RGBColor(25, 70, 98)

    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for i, b in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = Pt(22)


if __name__ == "__main__":
    prs = Presentation()
    for title, bullets in slides:
        add_slide(prs, title, bullets)
    prs.save("docs/presales/屯象OS-徐记海鲜-商务售前方案.pptx")
    print("Generated: docs/presales/屯象OS-徐记海鲜-商务售前方案.pptx")
